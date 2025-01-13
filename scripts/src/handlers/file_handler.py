import os
import shutil
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
from azure.storage.blob import BlobServiceClient
from io import BytesIO


class FileHandler:

    def __init__(self, temp_folder_path):
        self.connection_string = os.getenv('AZURE_STORAGE_CONNECTION_STRING')
        self.container_name = os.getenv('AZURE_BLOB_CONTAINER_NAME')
        self.blob_name_pptx_template = os.getenv('AZURE_BLOB_FILE_NAME_PPTX_TEMPLATE') 
        self.container_name_iuc_backup = os.getenv('AZURE_BLOB_CONTAINER_NAME_IUC_BACKUP')

        # Validate environment variables
        self._validate_env_variables()

        self.pptx_template_path = self.get_blob_pptx_template()
        self.screenshots_path = temp_folder_path
        self.new_file_path = None
        self.file_name_prefix = "IUC"
    
    def _validate_env_variables(self):
        """Validate that all the necessary environment variables are set."""
        if not self.connection_string:
            raise ValueError("Connection string not found in environment variables!")
        if not self.container_name:
            raise ValueError("Container name not found in environment variables!")
        if not self.blob_name_pptx_template:
            raise ValueError("Blob file name for PPTX template not found in environment variables!")
        if not self.container_name_iuc_backup:
            raise ValueError("Container name not found in environment variables!")

    def get_blob_pptx_template(self):
        try:
            # Get the current date and time
            current_datetime = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            print(f"Using container: {self.container_name}, blob: {self.blob_name_pptx_template}, at {current_datetime}")

            # 3. Initialize BlobServiceClient using the connection string
            blob_service_client = BlobServiceClient.from_connection_string(self.connection_string)
            container_client = blob_service_client.get_container_client(self.container_name)

            # 4. Download the PowerPoint file from Azure Blob Storage
            blob_client = container_client.get_blob_client(self.blob_name_pptx_template)
            blob_data = blob_client.download_blob()

            # 5. Read the downloaded data into memory (in-memory file)
            file_bytes = blob_data.readall()  # Read the file content into bytes
            pptx_data = BytesIO(file_bytes)  # Convert byte content to a BytesIO object (in-memory file)

            return pptx_data

        except ValueError as ve:
            print(f"ValueError occurred: {ve}")
        except Exception as e:
            print(f"An error occurred: {e}")


    def add_screenshots_to_template(self):
        try:
            # Creating a unique file path
            self.new_file_path = self.generate_file_path()

            # Load the PowerPoint template
            prs = Presentation(self.pptx_template_path)

            # Folder containing the screenshots
            screenshots_path = os.path.abspath(self.screenshots_path)

            # Ensure the folder exists
            if not os.path.exists(screenshots_path):
                print(f"The directory {self.screenshots_path} does not exist.")
                return

            # Get a list of all screenshot files in the folder
            screenshot_files = [
                f
                for f in os.listdir(screenshots_path)
                if f.endswith((".png", ".jpg", ".jpeg"))
            ]

            if not screenshot_files:
                print(f"No images found in {self.screenshots_path}.")
                return

            # Loop through the screenshots and add them to the presentation
            for i in range(0, len(screenshot_files), 2):
                # Choose the blank slide layout
                layout = prs.slide_layouts[27]  # Blank slide layout
                slide = prs.slides.add_slide(layout)

                # Define the paths to the images
                screenshot1 = screenshot_files[i]
                image_path1 = os.path.join(screenshots_path, screenshot1)

                # Place the first image (top image)
                slide.shapes.add_picture(
                    image_path1,
                    Inches(0.2),
                    Inches(1),
                    width=Inches(7),
                    height=Inches(4),
                )

                # Ensure there is a second image
                if i + 1 < len(screenshot_files):
                    screenshot2 = screenshot_files[i + 1]
                    image_path2 = os.path.join(screenshots_path, screenshot2)

                    # Place the second image (bottom image)
                    slide.shapes.add_picture(
                        image_path2,
                        Inches(0.2),
                        Inches(5.2),
                        width=Inches(7),
                        height=Inches(4),
                    )

            # Save the modified PowerPoint presentation
            prs.save(self.new_file_path)
            print(f"PowerPoint presentation saved as {self.new_file_path}")

            # Remove the temporary screenshots folder
            self.remove_temporary_folder()

            # Upload the new PowerPoint file to Azure Blob Storage
            self.upload_pptx_to_azure()

            return self.new_file_path

        except FileNotFoundError as e:
            print(f"Error: The file was not found. {e}")
        except PermissionError as e:
            print(f"Error: Permission denied. {e}")
        except Exception as e:
            print(f"An unexpected error occurred: {e}")

    def generate_file_path(self):
        try:
            # Use the generate_unique_filename method
            unique_filename = self.generate_unique_filename()

            if unique_filename:
                # Get the current working directory (cross-platform)
                current_directory = Path(__file__).parent.parent.parent

                # Define the relative base directory for your "IUCs" folder
                base_directory = current_directory / 'IUCs'
                
                # Construct the full file path using pathlib (handles different OS path formats)
                file_path = base_directory / f"{unique_filename}.pptx"

                print(f"Generated file path: {file_path}")

                # Convert to string to return it as a file path
                return str(file_path)
            else:
                return None
        except Exception as e:
            # Catch any exception
            print(f"Error occurred while generating the file path: {e}")
            return None
        

    def generate_unique_filename(self):
        try:
            # Get the current timestamp in a specific format (includes hour, minute, and second)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            # Combine the prefix and timestamp to generate a unique name
            unique_filename = f"{self.file_name_prefix}_{timestamp}"
            return unique_filename
        except Exception as e:
            # Catch any exception
            print(f"Error occurred while generating the unique name: {e}")
            return None
        
    def remove_file(self, file_path):
        """Remove the file after the email has been sent."""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                print(f"Successfully removed the file: {file_path}")
            else:
                print(f"File not found: {file_path}")
        except Exception as e:
            print(f"Error removing file: {e}")

    def remove_temporary_folder(self):
        try:
            # Get the full path of the temp folder (based on self.screenshots_path)
            screenshot_folder_to_remove = self.screenshots_path
            
            # Check if the folder exists and is a directory
            if os.path.exists(screenshot_folder_to_remove) and os.path.isdir(screenshot_folder_to_remove):
                # Remove the folder and all its contents
                shutil.rmtree(screenshot_folder_to_remove)
                print(f"Temporary folder '{screenshot_folder_to_remove}' and its contents cleared.")
            else:
                print(f"Error: The folder '{screenshot_folder_to_remove}' does not exist or is not a valid directory.")
        
        except FileNotFoundError as e:
            print(f"Error: The file was not found. {e}")
        except PermissionError as e:
            print(f"Error: Permission denied. {e}")
        except Exception as e:
            print(f"An unexpected error occurred: {e}")

    def upload_pptx_to_azure(self):

        """Upload the PowerPoint file to Azure Blob Storage."""
        try:
            # Initialize the BlobServiceClient
            blob_service_client = BlobServiceClient.from_connection_string(self.connection_string)
            blob_client = blob_service_client.get_blob_client(self.container_name_iuc_backup, self.new_file_path)

            # Upload the file to Azure Blob Storage
            with open(self.new_file_path, "rb") as data:
                blob_client.upload_blob(data, overwrite=True)

            print(f"Uploaded {self.new_file_path} to Azure Blob Storage.")
        except Exception as e:
            print(f"Error uploading to Azure: {e}")